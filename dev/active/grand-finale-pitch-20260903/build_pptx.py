from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "deck.pptx"
ASSETS = ROOT / "assets" / "live"

BG = RGBColor(1, 1, 2)
SURFACE = RGBColor(15, 16, 17)
HAIR = RGBColor(35, 37, 42)
INK = RGBColor(247, 248, 248)
MUTED = RGBColor(176, 181, 190)
SUBTLE = RGBColor(120, 125, 135)
ACCENT = RGBColor(94, 106, 210)
GREEN = RGBColor(39, 166, 68)
FONT = "Liberation Sans"
MONO = "Liberation Mono"

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=BG, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=24, color=INK, bold=False,
         font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0,
         name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 0.95 if size >= 42 else 1.05
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def line(slide, x1, y1, x2, y2, color=HAIR, width=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(max(y2-y1, 0.01)))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = color
    s.height = Pt(width)
    return s


def header(slide, right=""):
    rect(slide, 0, 0, 13.333, 7.5, BG)
    rect(slide, .36, .32, .09, .09, ACCENT, ACCENT, True)
    text(slide, "BIDPILOT", .51, .25, 2, .3, 10, MUTED, False, FONT)
    if right:
        text(slide, right.upper(), 8.7, .25, 4.25, .3, 10, MUTED, False, FONT, PP_ALIGN.RIGHT)
    line(slide, .36, .61, 12.97, .61)


def footer(slide, label, number):
    line(slide, .36, 7.08, 12.97, 7.08)
    text(slide, label, .36, 7.14, 7.5, .2, 8, SUBTLE)
    text(slide, f"{number:02d} / 10", 11.8, 7.14, 1.15, .2, 9, INK, False, MONO, PP_ALIGN.RIGHT)
    for i in range(10):
        color = ACCENT if i == number - 1 else HAIR
        rect(slide, .36 + i * 1.26, 7.43, 1.20, .02, color, color)


def title(slide, value, y=1.25, size=42, w=11.8):
    return text(slide, value, .36, y, w, 1.2, size, INK, True, FONT, name="Editable title")


def add_image_cover(slide, path, x, y, w, h, crop=(0, 0, 0, 0)):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    frame_ratio = w / h
    image_ratio = iw / ih
    if image_ratio > frame_ratio:
        shown = iw * frame_ratio / image_ratio
        left = (iw - shown) / 2 / iw
        right = left
        top = bottom = 0
    else:
        shown = ih * image_ratio / frame_ratio
        top = (ih - shown) / 2 / ih
        bottom = top
        left = right = 0
    left += crop[0]
    top += crop[1]
    right += crop[2]
    bottom += crop[3]
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.crop_left = max(0, left)
    pic.crop_top = max(0, top)
    pic.crop_right = max(0, right)
    pic.crop_bottom = max(0, bottom)
    return pic


def rubric(slide, percent, label):
    rect(slide, 9.85, .31, .05, .05, ACCENT, ACCENT, True)
    text(slide, f"{percent}%", 9.98, .24, .55, .25, 10, INK, True)
    text(slide, label.upper(), 10.54, .24, 2.42, .25, 10, MUTED, False, FONT, PP_ALIGN.RIGHT)


def numbered_rows(slide, rows, x, y, w, row_h=.52, size=18):
    for idx, value in enumerate(rows, start=1):
        yy = y + (idx - 1) * row_h
        line(slide, x, yy + row_h - .05, x + w, yy + row_h - .05)
        text(slide, f"{idx:02d}", x, yy + .04, .35, .28, 9, SUBTLE, False, MONO)
        text(slide, value, x + .42, yy, w - .42, .36, size, INK)


# Slide 1
s = prs.slides.add_slide(blank)
header(s, "Ten minute finale")
text(s, "SNOWFLAKE COCO CLI HACKATHON 2026 · GRAND FINALE", .36, 1.45, 7.0, .3, 13, SUBTLE)
title(s, "BidPilot", 2.05, 72, 7.0)
rect(s, .36, 3.45, .58, .02, ACCENT, ACCENT)
text(s, "Win the score, not the prompt.", .36, 3.68, 6.2, .6, 25, MUTED)
text(s, "Sergio Lee · Washington State CPA · sergiobuilds", .36, 6.78, 6.5, .2, 10, SUBTLE)
footer(s, "Grand Finale", 1)

# Slide 2
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 30, "Real-world relevance")
title(s, "The bid is often lost\nbefore writing begins.", 1.08, 40, 8.1)
labels = [("01", "Eligibility", "Can this supplier credibly bid."), ("02", "Score", "Where the official weight sits."), ("03", "Evidence", "What the supplier can prove."), ("04", "Ownership", "Who is accountable for the work.")]
for i, (n, a, b) in enumerate(labels):
    x = .36 + i * 3.16
    line(s, x, 3.45, x + 2.95, 3.45)
    text(s, n, x, 3.63, .4, .2, 9, ACCENT, False, MONO)
    text(s, a, x, 3.95, 2.8, .3, 20, INK)
    text(s, b, x, 4.42, 2.65, .65, 12, SUBTLE)
text(s, "Proposal teams commit resources before the buyer’s score and the supplier’s proof agree.", .36, 5.55, 10.6, .5, 16, SUBTLE)
footer(s, "The problem", 2)

# Slide 3
s = prs.slides.add_slide(blank)
header(s, "The contract")
title(s, "One accountable chain.", 1.27, 43)
chain = ["public tender + supplier evidence", "PURSUE · REVIEW · NO-GO", "score-weighted Win Position", "evidence-bound proposal", "red-team", "owned work", "same-run Snowflake readback"]
positions = [(.36, 2.65, 2.46), (3.05, 2.65, 2.25), (.72, 3.37, 2.54), (3.52, 3.37, 2.17), (5.95, 3.37, 1.15), (2.0, 4.09, 1.35), (3.62, 4.09, 2.55)]
for i, (value, pos) in enumerate(zip(chain, positions)):
    x, y, w = pos
    rect(s, x, y, w, .46, SURFACE, HAIR, True)
    text(s, value, x + .12, y + .10, w - .24, .22, 12, INK, False, FONT, PP_ALIGN.CENTER)
    if i in (0, 1, 2, 3, 4, 5):
        text(s, "→", x + w + .05, y + .11, .25, .2, 12, SUBTLE)
text(s, "A general LLM summarizes. BidPilot preserves the decision and execution contract.", .36, 5.05, 8.7, .35, 15, SUBTLE)
footer(s, "What the product is", 3)

# Slide 4
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 30, "Real-world relevance")
title(s, "REVIEW is a\ntrusted answer.", 1.10, 35, 4.0)
numbered_rows(s, ["Real public tender", "Synthetic demo supplier", "4 evidence gaps", "No run created"], .36, 3.2, 3.65, .58, 16)
text(s, "LIVE PRODUCT · TENDER DETAIL", 5.0, .83, 4.1, .25, 10, SUBTLE, False, MONO)
add_image_cover(s, ASSETS / "tender.png", 5.0, 1.13, 7.95, 5.55)
footer(s, "Live demonstration · real source", 4)

# Slide 5
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 30, "Solution completeness")
title(s, "Separate verified replay.", .93, 35, 7.0)
text(s, "SEPARATE SYNTHETIC HISTORICAL REPLAY", 8.0, 1.04, 4.9, .25, 12, SUBTLE, False, MONO, PP_ALIGN.RIGHT)
add_image_cover(s, ASSETS / "walkthrough.png", .36, 2.05, 12.60, 4.15)
for i, item in enumerate(["PURSUE", "40 points", "3 compared", "1 selected"]):
    text(s, item, .36 + i * 1.55, 6.35, 1.45, .3, 16, INK)
footer(s, "Live demonstration · verified replay", 5)

# Slide 6
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 30, "Solution completeness")
title(s, "The score controls\nthe work.", .92, 34, 5.2)
facts = ["4 weighted plans", "8 proposal sections", "Red-team passed", "12 owned and review tasks"]
numbered_rows(s, facts, 8.8, .92, 4.15, .45, 14)
text(s, "REPLAY · PROPOSAL AND RED-TEAM", .36, 2.65, 4, .2, 9, SUBTLE, False, MONO)
add_image_cover(s, ASSETS / "replay-03.png", .36, 2.93, 12.60, 1.65)
text(s, "REPLAY · OWNED WORK", .36, 4.77, 4, .2, 9, SUBTLE, False, MONO)
add_image_cover(s, ASSETS / "replay-04.png", .36, 5.05, 12.60, 1.55)
footer(s, "Live demonstration · execution readiness", 6)

# Slide 7
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 40, "Technical execution")
title(s, "Snowflake is the\noperating memory.", 1.05, 39, 7.0)
arch = [("01 · STORE", "Opportunity\nGraph"), ("02 · GOVERN", "Snowpark policy"), ("03 · RUN", "Cortex Code CLI"), ("04 · READ", "Streamlit reader")]
for i, (n, label) in enumerate(arch):
    x = .36 + i * 3.15
    line(s, x, 3.6, x + 2.95, 3.6)
    text(s, n, x, 3.78, 2.7, .2, 9, ACCENT, False, MONO)
    text(s, label, x, 4.15, 2.7, .55, 19, INK)
notes = ["Versioned tender + governed supplier evidence", "Runner writes. Reader replays complete runs.", "Authenticated failures fail closed."]
for i, note in enumerate(notes):
    text(s, note, .36 + i * 4.2, 5.4, 3.85, .7, 13, MUTED)
footer(s, "Architecture", 7)

# Slide 8
s = prs.slides.add_slide(blank)
header(s)
rubric(s, 40, "Technical execution")
title(s, "One recorded run.", 1.58, 42)
nums = [("1", "decision"), ("3", "strategies"), ("4", "plans"), ("8", "sections"), ("12", "tasks")]
for i, (num, label) in enumerate(nums):
    x = .36 + i * 2.53
    line(s, x, 3.25, x + 2.30, 3.25)
    text(s, num, x, 3.43, 2.0, .7, 42, INK, True)
    text(s, label, x, 4.27, 2.0, .3, 12, SUBTLE)
text(s, "cortex-final-20260802-a", .36, 5.18, 2.55, .25, 11, RGBColor(130, 140, 255), False, MONO)
text(s, "Cortex session + Snowflake query provenance · same run_id", 2.96, 5.18, 7.5, .25, 11, SUBTLE, False, MONO)
footer(s, "Provenance", 8)

# Slide 9
s = prs.slides.add_slide(blank)
header(s, "Safety and buyer")
title(s, "Human approval stays\nin the loop.", 1.38, 40, 8.2)
steps = ["Source review", "Supplier evidence", "Pricing", "Final edit", "Legal submission"]
x = .36
for idx, step in enumerate(steps):
    w = [1.52, 1.86, 1.05, 1.18, 1.72][idx]
    rect(s, x, 3.65, w, .47, SURFACE, HAIR, True)
    text(s, step, x + .10, 3.78, w - .2, .2, 12, INK, False, FONT, PP_ALIGN.CENTER)
    x += w + .32
    if idx < len(steps) - 1:
        text(s, "→", x - .24, 3.78, .18, .2, 11, SUBTLE)
text(s, "BidPilot prepares and governs the work. The first buyer is a small B2G proposal team.", .36, 4.65, 10.3, .35, 15, SUBTLE)
text(s, "Commercial path: team subscription + usage-based runs.", .36, 5.18, 8.4, .35, 15, MUTED)
footer(s, "Human boundary", 9)

# Slide 10
s = prs.slides.add_slide(blank)
header(s, "Close")
title(s, "Win the score,\nnot the prompt.", 1.02, 55, 8.0)
rect(s, .36, 3.43, .88, .02, ACCENT, ACCENT)
numbered_rows(s, ["A defensible pursuit decision.", "A score-weighted win strategy.", "Execution-ready proposal work."], .36, 3.83, 7.3, .58, 17)
text(s, "Persisted and replayable in Snowflake.", .36, 6.55, 5.5, .25, 12, SUBTLE)
footer(s, "BidPilot", 10)

for slide in prs.slides:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

prs.save(OUT)
print(OUT)
